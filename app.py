import os
import time
import random
from io import BytesIO
import glob
import streamlit as st
import pandas as pd
import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
from anthropic import Anthropic, APIStatusError
from PyPDF2 import PdfReader
from dotenv import load_dotenv

load_dotenv()

# ---------------------------
# API Key
# ---------------------------
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
if not ANTHROPIC_API_KEY:
    try:
        ANTHROPIC_API_KEY = st.secrets["ANTHROPIC_API_KEY"]
    except Exception:
        st.error("❌ No Anthropic API key found. Set it in environment variables or secrets.toml")
        st.stop()

client = Anthropic(api_key=ANTHROPIC_API_KEY)

# ---------------------------
# Embeddings (cached + safe)
# ---------------------------
@st.cache_resource
def load_embedder():
    try:
        return SentenceTransformer("all-MiniLM-L6-v2", device="cpu")
    except Exception:
        return None

embedder = load_embedder()
if embedder is None:
    st.error("⚠️ Could not load embedding model. Please reload the app.")
    if st.button("🔄 Reload App"):
        st.rerun()
    st.stop()

# ---------------------------
# Helpers
# ---------------------------
def chunk_text(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunks.append(" ".join(words[i:i + chunk_size]))
    return chunks

@st.cache_data(show_spinner=False)
def extract_text_from_file(file_bytes: bytes, file_name: str):
    ext = file_name.split(".")[-1].lower()
    try:
        if ext == "pdf":
            reader = PdfReader(BytesIO(file_bytes))
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            return text if text.strip() else None
        elif ext == "docx":
            doc = docx.Document(BytesIO(file_bytes))
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext in ["xlsx", "csv"]:
            df = pd.read_excel(BytesIO(file_bytes)) if ext=="xlsx" else pd.read_csv(BytesIO(file_bytes))
            return df.to_string()
        elif ext == "pptx":
            prs = Presentation(BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        elif ext == "txt":
            return file_bytes.decode("utf-8")
    except Exception:
        return None
    return None

@st.cache_resource(show_spinner=False)
def build_faiss_index(text: str):
    chunks = chunk_text(text)
    embeddings = embedder.encode(chunks, convert_to_numpy=True)
    dim = embeddings.shape[1]
    index = faiss.IndexFlatL2(dim)
    index.add(embeddings)
    return chunks, index

def search_chunks(query, chunks, index, top_k=3):
    try:
        query_emb = embedder.encode([query], convert_to_numpy=True)
        distances, indices = index.search(query_emb, top_k)
        results = []
        for idx in indices[0]:
            if idx >= 0 and idx < len(chunks):
                results.append(chunks[idx])
        return results
    except Exception:
        return []

def query_claude(client, model, messages, max_tokens=1000, retries=3, system=None):
    for attempt in range(retries):
        try:
            response = client.messages.create(
                model=model,
                max_tokens=max_tokens,
                system=system,
                messages=messages,
            )
            return response.content[0].text
        except APIStatusError as e:
            if e.status_code == 529:
                wait_time = (2 ** attempt) + random.uniform(0, 1)
                st.warning(f"⚠️ Claude API overloaded. Retrying in {wait_time:.1f}s...")
                time.sleep(wait_time)
                continue
            else:
                st.error("⚠️ API error. Please try again later.")
                return None
        except Exception:
            st.warning("⚠️ Something went wrong while contacting Claude. Try again or reload the app.")
            return None
    st.warning("🚨 Claude API is still overloaded. Please try again later.")
    return None

# ---------------------------
# Batch Summarizer
# ---------------------------
def batch_summarize(text, batch_size=5, use_web=False):
    chunks = chunk_text(text)
    summaries = []
    for i in range(0, len(chunks), batch_size):
        batch = "\n\n".join(chunks[i:i + batch_size])
        if use_web:
            prompt = (
                f"Summarize this section based on the document:\n\n{batch}\n\n"
                "Instructions: Use the document to create the summary. "
                "If the document lacks relevant information, you may add general knowledge. "
                "Do not omit document content if present."
            )
        else:
            prompt = f"Summarize this section strictly based on the text:\n\n{batch}"
        response = query_claude(
            client,
            "claude-opus-4-1-20250805",
            messages=[{"role": "user", "content": prompt}],
            system="You are an AI summarizer."
        )
        if response:
            summaries.append(response)
        time.sleep(0.5)
    combined = "\n".join(summaries)
    if use_web:
        final_prompt = (
            f"Combine these summaries into a single coherent summary using the document:\n\n{combined}\n\n"
            "Instructions: Ensure all document information is included. "
            "If the document lacks some info, you may append general knowledge. "
            "Do not remove or ignore document content."
        )
        system_msg = "Merge and enhance summaries with optional general knowledge. Document must always be included if available."
    else:
        final_prompt = f"Combine these summaries into a single coherent summary strictly based on the text:\n\n{combined}"
        system_msg = "Merge summaries strictly based on provided text."
    final_summary = query_claude(
        client,
        "claude-opus-4-1-20250805",
        messages=[{"role": "user", "content": final_prompt}],
        system=system_msg
    )
    return final_summary or "⚠️ Summarization failed. Please try again."

# ---------------------------
# Streamlit UI
# ---------------------------
st.set_page_config(page_title="AI Document Summarizer & NHAI Tenders", layout="wide")

st.markdown("""
<style>
body, .stApp, .stTextInput>div>div>input { font-size: 17px; }
.compact-label { font-weight: 500; margin-top: 0px; margin-bottom: 4px; }
.output-box { background-color: #f9f9f9; padding: 15px; border-radius: 10px; margin-top: 16px; font-size: 16px; line-height: 1.5; }
.output-box h1, .output-box h2, .output-box h3 { font-size: 18px !important; font-weight: 600; margin-top: 8px; margin-bottom: 6px; }
.output-box p { margin: 4px 0; font-size: 16px; }
.stSelectbox>div>div>div, .stTextInput>div { margin-top: 0px !important; margin-bottom: 4px !important; }
.stButton>button { width: 100%; }
</style>
""", unsafe_allow_html=True)

st.markdown("""<style>h1 { margin-top: -70px !important; margin-bottom: 12px !important; }</style>""", unsafe_allow_html=True)


st.title("📄 Document & Tender AI Assistant – Powered by Claude & RAG")

# ---------- Top-Level Mode Selector ----------
st.markdown('<div class="compact-label">Select Application Mode</div>', unsafe_allow_html=True)
app_mode = st.selectbox(
    "",
    ["Document Summarizer", "NHAI Tender Assistant"],
    label_visibility="collapsed"
)

col_sidebar, col_main = st.columns([1,3])

# =============================
# Document Summarizer Module
# =============================
if app_mode=="Document Summarizer":
    with col_sidebar:
        st.markdown('<div class="compact-label">Upload Document</div>', unsafe_allow_html=True)
        uploaded_file = st.file_uploader(
            "file_uploader",
            type=["pdf","docx","xlsx","csv","pptx","txt"],
            label_visibility="collapsed"
        )
        text_extracted = False
        if uploaded_file:
            file_bytes = uploaded_file.read()
            text = extract_text_from_file(file_bytes, uploaded_file.name)
            if text:
                text_extracted = True
                chunks, index = build_faiss_index(text)
                st.session_state.chunks = chunks
                st.session_state.index = index
                st.session_state.full_text = text
                ext = uploaded_file.name.split(".")[-1].lower()
                messages = {
                    "pdf": "✅ PDF processed successfully!",
                    "docx": "✅ Word document processed successfully!",
                    "xlsx": "✅ Spreadsheet processed successfully!",
                    "csv": "✅ Spreadsheet processed successfully!",
                    "pptx": "✅ Presentation processed successfully!",
                    "txt": "✅ Text file processed successfully!"
                }
                st.success(messages.get(ext, "✅ Document processed successfully!"))
            else:
                st.warning("⚠️ OCR is disabled on Streamlit Cloud. Please upload documents with extractable text.")

        if text_extracted:
            st.markdown('<div class="compact-label">Select Mode</div>', unsafe_allow_html=True)
            selected_mode = st.selectbox("mode_selector", ["💬 Ask a Question","📝 Get Summary"], key="selected_mode", label_visibility="collapsed")
            use_web = st.checkbox("🌐 Include general knowledge if needed", value=False, help="If enabled, answers and summaries may use Claude's general knowledge.")

            # ---------- Collapsible History ----------
            with st.expander("🕘 History"):
                if "qa_history" not in st.session_state:
                    st.session_state.qa_history = []
                if "summary_history" not in st.session_state:
                    st.session_state.summary_history = []

                with st.expander("🗨️ Q&A History"):
                    if st.session_state.qa_history:
                        for entry in reversed(st.session_state.qa_history[-5:]):
                            st.markdown(f'<div class="output-box"><b>Q:</b> {entry["query"]}<br><b>A:</b> {entry["answer"]}</div>', unsafe_allow_html=True)
                    else:
                        st.info("No Q&A history yet.")

                with st.expander("📝 Summary History"):
                    if st.session_state.summary_history:
                        for i, s in enumerate(reversed(st.session_state.summary_history[-5:]), 1):
                            st.markdown(f'<div class="output-box"><b>Summary {i}:</b><br>{s}</div>', unsafe_allow_html=True)
                    else:
                        st.info("No summaries yet.")

    # ---------- Main ----------
    with col_main:
        if text_extracted:
            if selected_mode=="💬 Ask a Question":
                st.markdown('<div class="compact-label">Your Question</div>', unsafe_allow_html=True)
                user_query = st.text_input("question_input", placeholder="Type your question here...", label_visibility="collapsed")
                ask_disabled = user_query.strip()=="" 
                if st.button("Ask Now", disabled=ask_disabled):
                    with st.spinner(f"🧠 Searching for answers in your {ext}..."):
                        context_chunks = search_chunks(user_query, st.session_state.chunks, st.session_state.index)
                        context_text = "\n".join(context_chunks) if context_chunks else ""
                        if use_web:
                            prompt = f"Context:\n{context_text}\n\nQuestion: {user_query}\nInstructions: Answer using document first, then general knowledge if needed."
                            system_msg = "Answer using document and general knowledge if necessary."
                        else:
                            prompt = f"Context:\n{context_text}\n\nQuestion: {user_query}"
                            system_msg = "Answer strictly based on document."
                        response = query_claude(client, "claude-opus-4-1-20250805", messages=[{"role":"user","content":prompt}], system=system_msg)
                        if response:
                            st.session_state.qa_history.append({"query":user_query,"answer":response})
                            st.markdown(f'<div class="output-box"><b>🤖 Answer:</b><br>{response}</div>', unsafe_allow_html=True)
                        else:
                            st.warning("⚠️ Could not fetch an answer. Try again.")

            elif selected_mode=="📝 Get Summary":
                if st.button("Generate Summary"):
                    with st.spinner(f"📝 Generating a concise summary of your {ext}..."):
                        summary = batch_summarize(st.session_state.full_text, use_web=use_web)
                        st.session_state.summary_history.append(summary)
                        st.markdown(f'<div class="output-box"><b>📌 Document Summary:</b><br>{summary}</div>', unsafe_allow_html=True)


# =============================
# NHAI Tender Assistant Module (FAISS + RAG)
# =============================
if app_mode == "NHAI Tender Assistant":
    with col_sidebar:
        # Tender Source
        st.markdown('<div class="compact-label">Select Tender Source</div>', unsafe_allow_html=True)
        tender_source = st.selectbox(
            "tender_source_select",
            ["Use Existing Tender", "Upload New Tender"],
            label_visibility="collapsed"
        )

        selected_tender_name = None
        uploaded_tender_files = None
        tender_key = None

        # Initialize tender_history dict if not present
        if "tender_history" not in st.session_state:
            st.session_state.tender_history = {}

        # Select Existing Tender
        if tender_source == "Use Existing Tender":
            st.markdown('<div class="compact-label">Select Tender</div>', unsafe_allow_html=True)
            tender_folders = [f for f in os.listdir("tenders") if os.path.isdir(os.path.join("tenders", f))]
            selected_tender_name = st.selectbox(
                "Select Tender",
                ["-- Select Tender --"] + tender_folders,
                index=0,
                label_visibility="collapsed"
            )
            if selected_tender_name != "-- Select Tender --":
                tender_key = selected_tender_name
        else:
            # Upload New Tender
            st.markdown('<div class="compact-label">Upload Tender PDF(s)</div>', unsafe_allow_html=True)
            uploaded_tender_files = st.file_uploader(
                "",
                type=["pdf"],
                accept_multiple_files=True,
                label_visibility="collapsed"
            )
            if uploaded_tender_files:
                tender_key = "uploaded_tender"

        # Initialize session state for this tender_key if it exists
        if tender_key and tender_key not in st.session_state.tender_history:
            st.session_state.tender_history[tender_key] = {
                "full_text": "",
                "chunks": [],
                "index": None,
                "summary": None,
                "qa": []
            }

    # ----------------------------
    # Main Column
    # ----------------------------
    with col_main:
        full_text = ""
        text_extracted = False

        if tender_key:
            # Extract text from existing tenders
            if selected_tender_name and tender_source == "Use Existing Tender":
                folder_path = os.path.join("tenders", selected_tender_name)
                for fpath in glob.glob(os.path.join(folder_path, "*.pdf")):
                    try:
                        with open(fpath, "rb") as f:
                            content = extract_text_from_file(f.read(), fpath)
                            if content and content.strip():
                                full_text += "\n\n" + content
                                text_extracted = True
                    except Exception:
                        pass

            # Extract text from uploaded PDFs
            if uploaded_tender_files:
                for file in uploaded_tender_files:
                    try:
                        content = extract_text_from_file(file.read(), file.name)
                        if content and content.strip():
                            full_text += "\n\n" + content
                            text_extracted = True
                    except Exception:
                        pass

            # Show OCR warning ONLY if user tried to provide a tender but no text found
            if not text_extracted:
                st.warning("⚠️ OCR is disabled on Streamlit Cloud. Please upload documents with extractable text.")
                st.stop()

            # Save text in session state
            st.session_state.tender_history[tender_key]["full_text"] = full_text

            # Build FAISS index only once
            if not st.session_state.tender_history[tender_key]["index"]:
                chunks, index = build_faiss_index(full_text)
                st.session_state.tender_history[tender_key]["chunks"] = chunks
                st.session_state.tender_history[tender_key]["index"] = index

            # Generate summary only if not already done
            if not st.session_state.tender_history[tender_key]["summary"]:
                with st.spinner("📝 Generating tender summary..."):
                    prompt = f"Extract key tender details and create a professional summary:\n\n{full_text[:20000]}"
                    summary = query_claude(
                        client,
                        "claude-opus-4-1-20250805",
                        messages=[{"role": "user", "content": prompt}],
                        system="You are an expert summarizer for NHAI tenders."
                    )
                    st.session_state.tender_history[tender_key]["summary"] = summary

            # Display summary
            st.markdown('<div class="tender-section-title">📌 Key Tender Information & Summary</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="output-box">{st.session_state.tender_history[tender_key]["summary"]}</div>', unsafe_allow_html=True)

            # ----------------------------
            # Ask questions ONLY if summary is ready
            # ----------------------------
            st.markdown('<div class="tender-qa-title">💬 Ask Question</div>', unsafe_allow_html=True)
            tender_query = st.text_input(
                "tender_question_input",
                placeholder="Type your question here...",
                label_visibility="collapsed"
            )

            if st.button("Ask Question", key="ask_tender_btn") and tender_query.strip():
                # Check history first
                existing = [q for q in st.session_state.tender_history[tender_key]["qa"] if q["query"] == tender_query]
                if existing:
                    answer = existing[0]["answer"]
                else:
                    chunks = st.session_state.tender_history[tender_key]["chunks"]
                    index = st.session_state.tender_history[tender_key]["index"]
                    relevant_chunks = search_chunks(tender_query, chunks, index)
                    context_text = "\n".join(relevant_chunks) if relevant_chunks else full_text[:2000]

                    prompt = f"Answer the following question strictly based on the tender document:\n\n{context_text}\n\nQuestion:\n{tender_query}"
                    answer = query_claude(
                        client,
                        "claude-opus-4-1-20250805",
                        messages=[{"role": "user", "content": prompt}],
                        system="You are an expert answering questions from tender documents."
                    )
                    st.session_state.tender_history[tender_key]["qa"].append({"query": tender_query, "answer": answer})

                st.markdown(f'<div class="output-box"><b>🤖 Answer:</b> {answer}</div>', unsafe_allow_html=True)



# ---------- Fixed Footer Disclaimer ----------
st.markdown(
    """
    <div style="
        position: fixed;
        bottom: 0;
        left: 0;
        width: 100%;
        background-color:#f0f0f0;
        padding:10px;
        border-top: 1px solid #ccc;
        font-size:12px;
        color:#333;
        text-align:center;
        z-index:1000;
    ">
    ⚠️ <b>Disclaimer:</b> This AI-powered document summarizer and chatbot is built using publicly available resources. 
    While it strives to provide accurate and helpful responses, it may occasionally make errors or omissions. 
    Please verify any critical information independently.
    </div>
    """,
    unsafe_allow_html=True
)
